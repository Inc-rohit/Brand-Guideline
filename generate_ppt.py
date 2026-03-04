"""
Incisiv Brand-Compliant PowerPoint Generator
Topic: Agentic AI in Retail
Brand tokens sourced from index.html design system
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ── Brand Tokens ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x13, 0x1c, 0x4e)
ORANGE     = RGBColor(0xFA, 0x96, 0x28)
PURPLE     = RGBColor(0x7B, 0x2D, 0x8B)
TEAL       = RGBColor(0x14, 0x31, 0x4D)
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT       = RGBColor(0x11, 0x11, 0x11)
TEXT2      = RGBColor(0x55, 0x55, 0x55)
TEXT3      = RGBColor(0x99, 0x99, 0x99)
PAGE_BG    = RGBColor(0xF7, 0xF7, 0xF5)
BORDER     = RGBColor(0xE5, 0xE5, 0xE5)

FONT_HEAD  = "Poppins"
FONT_BODY  = "Libre Franklin"

# ── Slide size: 16:9 widescreen ─────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font=FONT_HEAD, size=Pt(24), bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font=FONT_BODY, size=Pt(13), bold=False,
             color=TEXT2, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def add_bullet_box(slide, items, x, y, w, h,
                   font=FONT_BODY, size=Pt(13), color=TEXT2,
                   bullet_color=ORANGE):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"— {item}"
        run.font.name  = font
        run.font.size  = size
        run.font.color.rgb = color
    return txb


def add_label(slide, text, x, y, w, h):
    """Small all-caps orange label."""
    add_textbox(slide, text.upper(), x, y, w, h,
                font=FONT_HEAD, size=Pt(9), bold=True,
                color=ORANGE, align=PP_ALIGN.LEFT)


def orange_bar(slide, y=Inches(0.5), thickness=Pt(2.5)):
    """Thin orange horizontal accent bar at top."""
    add_rect(slide, 0, y, W, Emu(thickness), fill_rgb=ORANGE)


def footer(slide, label="Incisiv • Confidential"):
    add_textbox(slide, label,
                Inches(0.5), H - Inches(0.42), W - Inches(1), Inches(0.35),
                font=FONT_BODY, size=Pt(8), color=TEXT3, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, H - Inches(0.04), W, Inches(0.04), fill_rgb=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════════════════
def slide_cover(prs):
    sl = prs.slides.add_slide(BLANK)

    # Full navy background
    add_rect(sl, 0, 0, W, H, fill_rgb=NAVY)

    # Left orange accent strip
    add_rect(sl, 0, 0, Inches(0.25), H, fill_rgb=ORANGE)

    # Bottom purple fade strip
    add_rect(sl, 0, H - Inches(1.6), W, Inches(1.6), fill_rgb=TEAL)

    # Decorative circle motif (top right)
    for r, a in [(Inches(5), 0.04), (Inches(3.8), 0.06), (Inches(2.4), 0.08)]:
        c = sl.shapes.add_shape(9, W - r, -r//2, r*2, r*2)   # oval
        c.fill.background()
        c.line.color.rgb = WHITE
        c.line.width = Pt(0.75)
        import ctypes
        try:
            c.line.fill.fore_color.rgb = WHITE
        except Exception:
            pass
        # set opacity-like via alpha isn't easy; just use low-opacity white equivalent
        c.line.color.rgb = RGBColor(0x1e, 0x2d, 0x6e)

    # Chapter label
    add_label(sl, "Incisiv Research  •  2026", Inches(0.6), Inches(1.2), Inches(6), Inches(0.4))

    # Main headline
    add_textbox(sl, "Agentic AI\nin Retail",
                Inches(0.6), Inches(1.8), Inches(8), Inches(3.2),
                font=FONT_HEAD, size=Pt(56), bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    # Sub-headline
    add_textbox(sl,
                "How autonomous AI agents are reshaping retail operations,\ncustomer experience, and competitive strategy.",
                Inches(0.6), Inches(4.85), Inches(9), Inches(1.1),
                font=FONT_BODY, size=Pt(15), bold=False,
                color=RGBColor(0xB0, 0xBB, 0xD8), align=PP_ALIGN.LEFT)

    # Orange rule
    add_rect(sl, Inches(0.6), Inches(4.7), Inches(1.4), Emu(Pt(2.5)),
             fill_rgb=ORANGE)

    # Bottom meta
    add_textbox(sl, "Incisiv Intelligence Platform  |  incisiv.io",
                Inches(0.6), H - Inches(0.9), Inches(6), Inches(0.4),
                font=FONT_BODY, size=Pt(9), color=TEXT3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What Is Agentic AI?  (definition / intro)
# ════════════════════════════════════════════════════════════════════════════
def slide_definition(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "1.0  Setting the Stage", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))

    add_textbox(sl, "What Is Agentic AI?",
                Inches(0.7), Inches(0.75), Inches(9), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    # Definition pull-quote box (navy left border)
    add_rect(sl, Inches(0.7), Inches(1.8), Inches(0.06), Inches(1.4), fill_rgb=ORANGE)
    add_textbox(sl,
                '"Agentic AI systems are autonomous software agents that perceive their environment, '
                'reason over goals, plan multi-step actions, and execute tasks — with minimal human intervention."',
                Inches(0.95), Inches(1.75), Inches(8.8), Inches(1.5),
                font=FONT_BODY, size=Pt(14), italic=True, color=TEXT,
                align=PP_ALIGN.LEFT)

    # Three pillars
    pillars = [
        ("Perceive",  NAVY,   "Ingest real-time data from POS, inventory, CRM, web, and IoT sensors."),
        ("Reason",    PURPLE, "Use LLMs and planning algorithms to break goals into actionable steps."),
        ("Act",       ORANGE, "Execute decisions — place orders, adjust prices, respond to customers."),
    ]
    box_w = Inches(3.6)
    gap   = Inches(0.42)
    start_x = Inches(0.7)
    top   = Inches(3.7)

    for i, (title, col, desc) in enumerate(pillars):
        bx = start_x + i * (box_w + gap)
        add_rect(sl, bx, top, box_w, Inches(2.6),
                 fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(1))
        # color top bar
        add_rect(sl, bx, top, box_w, Inches(0.18), fill_rgb=col)
        add_textbox(sl, title, bx + Inches(0.28), top + Inches(0.32), box_w - Inches(0.4), Inches(0.5),
                    font=FONT_HEAD, size=Pt(18), bold=True, color=col)
        add_textbox(sl, desc, bx + Inches(0.28), top + Inches(0.95), box_w - Inches(0.4), Inches(1.5),
                    font=FONT_BODY, size=Pt(12), color=TEXT2)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Retail Opportunity
# ════════════════════════════════════════════════════════════════════════════
def slide_opportunity(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "2.0  The Opportunity", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "Why Retail? Why Now?",
                Inches(0.7), Inches(0.75), Inches(9), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    # Stat cards row
    stats = [
        ("$9.2T",   "Global retail market size in 2025"),
        ("73%",     "of retailers cite operational complexity as their top challenge"),
        ("3.4×",    "productivity lift reported by early agentic-AI adopters"),
        ("2027",    "Year most large retailers plan to deploy autonomous agents at scale"),
    ]
    card_w = Inches(2.9)
    gap = Inches(0.26)
    sx  = Inches(0.7)
    sy  = Inches(1.9)

    for i, (stat, label) in enumerate(stats):
        cx = sx + i * (card_w + gap)
        add_rect(sl, cx, sy, card_w, Inches(1.8),
                 fill_rgb=NAVY, line_rgb=None)
        add_textbox(sl, stat, cx + Inches(0.2), sy + Inches(0.18),
                    card_w - Inches(0.3), Inches(0.8),
                    font=FONT_HEAD, size=Pt(32), bold=True, color=ORANGE)
        add_textbox(sl, label, cx + Inches(0.2), sy + Inches(0.95),
                    card_w - Inches(0.3), Inches(0.75),
                    font=FONT_BODY, size=Pt(11), color=RGBColor(0xB0, 0xBB, 0xD8))

    # Paragraph context
    add_textbox(sl,
                "Retail is uniquely positioned to benefit from agentic AI: high transaction volumes, rich "
                "real-time data streams, razor-thin margins, and rapidly shifting consumer expectations "
                "create the perfect conditions for autonomous agents to generate measurable ROI.",
                Inches(0.7), Inches(4.0), Inches(12.0), Inches(1.0),
                font=FONT_BODY, size=Pt(13), color=TEXT2)

    # Key pressures list
    add_textbox(sl, "Pressures Accelerating Adoption",
                Inches(0.7), Inches(5.15), Inches(5), Inches(0.4),
                font=FONT_HEAD, size=Pt(13), bold=True, color=TEXT)
    add_bullet_box(sl,
                   ["Labour cost inflation and talent scarcity",
                    "Consumer expectations for hyper-personalisation",
                    "Supply chain volatility requiring real-time response",
                    "Competitive pressure from digitally-native disruptors"],
                   Inches(0.7), Inches(5.6), Inches(5.5), Inches(1.7),
                   size=Pt(12))

    # Right side image placeholder (branded)
    add_rect(sl, Inches(7.0), Inches(5.1), Inches(5.8), Inches(2.0),
             fill_rgb=TEAL, line_rgb=None)
    add_textbox(sl, "Retail AI Adoption is accelerating\nacross all formats and segments.",
                Inches(7.2), Inches(5.3), Inches(5.4), Inches(1.5),
                font=FONT_BODY, size=Pt(12), italic=True,
                color=RGBColor(0xB0, 0xBB, 0xD8))

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Use Cases
# ════════════════════════════════════════════════════════════════════════════
def slide_use_cases(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "3.0  Use Cases", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "Where Agents Are Deployed Today",
                Inches(0.7), Inches(0.75), Inches(11), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    use_cases = [
        ("Inventory & Replenishment",
         NAVY,
         ["Autonomous demand forecasting",
          "Auto-triggered purchase orders",
          "Out-of-stock prevention loops"]),
        ("Dynamic Pricing",
         PURPLE,
         ["Real-time competitor price monitoring",
          "Margin-aware markdown agents",
          "Promotion yield optimisation"]),
        ("Customer Experience",
         TEAL,
         ["Personalised product recommendations",
          "Autonomous chat & voice agents",
          "Post-purchase resolution bots"]),
        ("Supply Chain",
         RGBColor(0x1a, 0x6b, 0x5a),
         ["Supplier risk monitoring agents",
          "Logistics rerouting on disruption",
          "Shipment ETA communication"]),
        ("Store Operations",
         RGBColor(0x7a, 0x44, 0x00),
         ["Automated shelf-audit via vision AI",
          "Staff scheduling optimisation",
          "Energy & shrinkage management"]),
        ("Marketing & Campaigns",
         RGBColor(0x5c, 0x1a, 0x6e),
         ["Autonomous content generation",
          "Multi-channel campaign orchestration",
          "Real-time spend reallocation"]),
    ]

    cols = 3
    card_w = Inches(3.85)
    card_h = Inches(2.4)
    gap_x  = Inches(0.28)
    gap_y  = Inches(0.24)
    sx = Inches(0.7)
    sy = Inches(1.9)

    for i, (title, col, bullets) in enumerate(use_cases):
        row = i // cols
        c   = i %  cols
        cx  = sx + c * (card_w + gap_x)
        cy  = sy + row * (card_h + gap_y)

        add_rect(sl, cx, cy, card_w, card_h,
                 fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
        add_rect(sl, cx, cy, card_w, Inches(0.14), fill_rgb=col)

        add_textbox(sl, title,
                    cx + Inches(0.22), cy + Inches(0.22),
                    card_w - Inches(0.35), Inches(0.5),
                    font=FONT_HEAD, size=Pt(13), bold=True, color=col)

        add_bullet_box(sl, bullets,
                       cx + Inches(0.22), cy + Inches(0.82),
                       card_w - Inches(0.35), card_h - Inches(0.95),
                       size=Pt(11.5), color=TEXT2)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Adoption Maturity Curve
# ════════════════════════════════════════════════════════════════════════════
def slide_maturity(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "4.0  Maturity Landscape", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "Agentic AI Adoption Maturity Curve",
                Inches(0.7), Inches(0.75), Inches(10), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    stages = [
        ("Stage 1\nAssist",       "15%", BORDER,  TEXT3,
         "AI surfaces recommendations;\nhuman approves every action."),
        ("Stage 2\nAutomate",     "38%", RGBColor(0xd0,0xd8,0xf0), NAVY,
         "Routine tasks run autonomously;\nexceptions escalated to humans."),
        ("Stage 3\nOrchestrate",  "32%", NAVY,    WHITE,
         "Agent networks collaborate;\nhuman oversight at goal level."),
        ("Stage 4\nAutonomous",   "15%", ORANGE,  BLACK,
         "End-to-end agentic workflows;\nhuman sets objectives only."),
    ]

    bar_top  = Inches(2.1)
    bar_h    = Inches(1.8)
    sx       = Inches(0.7)
    total_w  = W - Inches(1.4)
    pcts     = [0.15, 0.38, 0.32, 0.15]

    x = sx
    for i, (label, pct, bg, fg, desc) in enumerate(stages):
        bw = Emu(int(total_w * pcts[i]))
        add_rect(sl, x, bar_top, bw, bar_h, fill_rgb=bg, line_rgb=None)
        add_textbox(sl, label, x + Inches(0.15), bar_top + Inches(0.22),
                    bw - Inches(0.2), Inches(0.75),
                    font=FONT_HEAD, size=Pt(11), bold=True, color=fg)
        add_textbox(sl, pct,
                    x + Inches(0.15), bar_top + Inches(1.0),
                    bw - Inches(0.2), Inches(0.5),
                    font=FONT_HEAD, size=Pt(20), bold=True, color=fg)
        x += bw

    add_textbox(sl, "% of retailers at each stage (Incisiv survey, 2025)",
                Inches(0.7), Inches(4.05), Inches(8), Inches(0.35),
                font=FONT_BODY, size=Pt(10), italic=True, color=TEXT3)

    # Stage descriptions below
    x = sx
    for i, (label, pct, bg, fg, desc) in enumerate(stages):
        bw = Emu(int(total_w * pcts[i]))
        add_textbox(sl, desc, x + Inches(0.1), Inches(4.5),
                    bw - Inches(0.15), Inches(1.2),
                    font=FONT_BODY, size=Pt(11), color=TEXT2)
        x += bw

    # Key insight box
    add_rect(sl, Inches(0.7), Inches(5.85), Inches(11.93), Inches(1.2),
             fill_rgb=NAVY, line_rgb=None)
    add_textbox(sl,
                "Key Insight:  Over half of retailers (53%) remain in assist/automate stages — "
                "signalling a significant window for orchestration-layer solutions and advisory.",
                Inches(0.95), Inches(5.98), Inches(11.5), Inches(0.95),
                font=FONT_BODY, size=Pt(13), color=WHITE)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Benefits
# ════════════════════════════════════════════════════════════════════════════
def slide_benefits(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "5.0  Business Impact", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "Measurable Benefits for Retailers",
                Inches(0.7), Inches(0.75), Inches(10), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    benefits = [
        ("↓ 22%",  "Inventory carrying\ncosts",          NAVY),
        ("↑ 18%",  "Conversion rate via\npersonalisation", PURPLE),
        ("↓ 35%",  "Customer service\nresolution time",   TEAL),
        ("↑ 12%",  "Gross margin via\ndynamic pricing",   RGBColor(0x1a, 0x6b, 0x5a)),
        ("↑ 40%",  "Supply chain\nresponsiveness",        RGBColor(0x7a, 0x44, 0x00)),
        ("↓ 28%",  "Labour hours on\nroutine tasks",       RGBColor(0x5c, 0x1a, 0x6e)),
    ]

    card_w = Inches(3.7)
    card_h = Inches(2.2)
    gap_x  = Inches(0.35)
    gap_y  = Inches(0.28)
    sx = Inches(0.7)
    sy = Inches(2.0)
    cols = 3

    for i, (metric, label, col) in enumerate(benefits):
        row = i // cols
        c   = i %  cols
        cx  = sx + c * (card_w + gap_x)
        cy  = sy + row * (card_h + gap_y)

        add_rect(sl, cx, cy, card_w, card_h, fill_rgb=col, line_rgb=None)
        add_textbox(sl, metric,
                    cx + Inches(0.25), cy + Inches(0.2),
                    card_w - Inches(0.4), Inches(0.85),
                    font=FONT_HEAD, size=Pt(30), bold=True, color=ORANGE)
        add_textbox(sl, label,
                    cx + Inches(0.25), cy + Inches(1.1),
                    card_w - Inches(0.4), Inches(0.9),
                    font=FONT_BODY, size=Pt(13), color=RGBColor(0xB0, 0xBB, 0xD8))

    add_textbox(sl, "Sources: Incisiv Agentic AI in Retail Benchmark Study, 2025. N=320 retailers.",
                Inches(0.7), H - Inches(0.65), Inches(10), Inches(0.3),
                font=FONT_BODY, size=Pt(9), italic=True, color=TEXT3)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Challenges & Considerations
# ════════════════════════════════════════════════════════════════════════════
def slide_challenges(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "6.0  Challenges", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "Key Challenges & Considerations",
                Inches(0.7), Inches(0.75), Inches(10), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    # Two-column: left = challenges, right = mitigations
    # Left header
    add_rect(sl, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.45), fill_rgb=NAVY)
    add_textbox(sl, "Challenges", Inches(0.9), Inches(1.88),
                Inches(5.2), Inches(0.4),
                font=FONT_HEAD, size=Pt(13), bold=True, color=WHITE)

    challenges = [
        ("Data Readiness",        "Siloed, inconsistent data limits agent accuracy and trust."),
        ("Governance & Control",  "Autonomous actions require clear guardrails, audit trails, and rollback."),
        ("Integration Complexity","Legacy systems rarely expose the APIs agents need to act."),
        ("Change Management",     "Workforce adoption and trust-building is as critical as the technology."),
        ("Hallucination Risk",    "LLM-powered agents can make plausible but incorrect decisions."),
    ]

    mitigations = [
        ("Invest in a Unified Data Layer",   "Retail data mesh or lakehouse as the agent's ground truth."),
        ("Define Agent Operating Boundaries","Human-in-the-loop thresholds and audit logging from day one."),
        ("API-First Modernisation",          "Prioritise integration middleware before scaling agents."),
        ("Upskilling & Co-pilot Phasing",    "Start with assist-mode agents; build confidence gradually."),
        ("Grounding & Retrieval Augmentation","RAG + structured data constraints to reduce hallucination."),
    ]

    add_rect(sl, Inches(6.9), Inches(1.85), Inches(5.9), Inches(0.45), fill_rgb=ORANGE)
    add_textbox(sl, "Mitigations", Inches(7.1), Inches(1.88),
                Inches(5.5), Inches(0.4),
                font=FONT_HEAD, size=Pt(13), bold=True, color=BLACK)

    row_h = Inches(0.82)
    sy    = Inches(2.4)

    for i, ((ch_title, ch_desc), (mi_title, mi_desc)) in enumerate(zip(challenges, mitigations)):
        cy = sy + i * row_h
        bg = WHITE if i % 2 == 0 else PAGE_BG

        add_rect(sl, Inches(0.7), cy, Inches(5.6), row_h,
                 fill_rgb=bg, line_rgb=BORDER, line_width=Pt(0.5))
        add_textbox(sl, ch_title, Inches(0.9), cy + Inches(0.06),
                    Inches(5.0), Inches(0.32),
                    font=FONT_HEAD, size=Pt(12), bold=True, color=NAVY)
        add_textbox(sl, ch_desc, Inches(0.9), cy + Inches(0.38),
                    Inches(5.1), Inches(0.35),
                    font=FONT_BODY, size=Pt(11), color=TEXT2)

        add_rect(sl, Inches(6.9), cy, Inches(5.9), row_h,
                 fill_rgb=bg, line_rgb=BORDER, line_width=Pt(0.5))
        add_textbox(sl, mi_title, Inches(7.1), cy + Inches(0.06),
                    Inches(5.5), Inches(0.32),
                    font=FONT_HEAD, size=Pt(12), bold=True, color=RGBColor(0x7a, 0x44, 0x00))
        add_textbox(sl, mi_desc, Inches(7.1), cy + Inches(0.38),
                    Inches(5.5), Inches(0.35),
                    font=FONT_BODY, size=Pt(11), color=TEXT2)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — The Road Ahead
# ════════════════════════════════════════════════════════════════════════════
def slide_roadahead(prs):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill_rgb=PAGE_BG)
    orange_bar(sl, y=0, thickness=Pt(5))

    add_label(sl, "7.0  Looking Ahead", Inches(0.7), Inches(0.35), Inches(6), Inches(0.4))
    add_textbox(sl, "The Road Ahead for Retail AI",
                Inches(0.7), Inches(0.75), Inches(10), Inches(0.9),
                font=FONT_HEAD, size=Pt(34), bold=True, color=TEXT)

    timeline = [
        ("2025",      ORANGE, "Foundation",
         ["Agent-assist tools widely deployed",
          "Inventory & pricing agents in production",
          "Pilot programmes for customer-facing agents"]),
        ("2026–27",   NAVY,   "Orchestration",
         ["Multi-agent frameworks in supply chain",
          "Personalisation agents at customer scale",
          "Retail-specific LLMs fine-tuned on private data"]),
        ("2028+",     PURPLE, "Autonomous Retail",
         ["End-to-end autonomous store operations",
          "Agent-to-agent B2B commerce",
          "Self-optimising retail enterprise"]),
    ]

    card_w = Inches(3.8)
    gap    = Inches(0.37)
    sx     = Inches(0.7)
    sy     = Inches(2.0)
    card_h = Inches(4.6)

    for i, (year, col, stage, bullets) in enumerate(timeline):
        cx = sx + i * (card_w + gap)
        add_rect(sl, cx, sy, card_w, card_h,
                 fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
        add_rect(sl, cx, sy, card_w, Inches(0.22), fill_rgb=col)

        add_textbox(sl, year, cx + Inches(0.25), sy + Inches(0.32),
                    card_w - Inches(0.4), Inches(0.6),
                    font=FONT_HEAD, size=Pt(26), bold=True, color=col)
        add_textbox(sl, stage, cx + Inches(0.25), sy + Inches(0.95),
                    card_w - Inches(0.4), Inches(0.4),
                    font=FONT_HEAD, size=Pt(14), bold=True, color=TEXT)

        add_rect(sl, cx + Inches(0.25), sy + Inches(1.42),
                 Inches(0.8), Emu(Pt(2)), fill_rgb=col)

        add_bullet_box(sl, bullets,
                       cx + Inches(0.25), sy + Inches(1.7),
                       card_w - Inches(0.4), card_h - Inches(1.9),
                       size=Pt(12.5), color=TEXT2)

    footer(sl)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Closing / Thank You
# ════════════════════════════════════════════════════════════════════════════
def slide_close(prs):
    sl = prs.slides.add_slide(BLANK)

    # Full navy background
    add_rect(sl, 0, 0, W, H, fill_rgb=NAVY)

    # Bottom orange rule
    add_rect(sl, 0, H - Inches(0.18), W, Inches(0.18), fill_rgb=ORANGE)

    # Left orange accent strip
    add_rect(sl, 0, 0, Inches(0.25), H, fill_rgb=ORANGE)

    add_textbox(sl, "Thank You",
                Inches(0.6), Inches(2.1), Inches(10), Inches(1.4),
                font=FONT_HEAD, size=Pt(52), bold=True, color=WHITE)

    add_rect(sl, Inches(0.6), Inches(3.6), Inches(1.6), Emu(Pt(2.5)), fill_rgb=ORANGE)

    add_textbox(sl,
                "To learn more about Incisiv's Agentic AI in Retail research,\n"
                "advisory services, and peer benchmarking programmes, contact us at:",
                Inches(0.6), Inches(3.85), Inches(9), Inches(1.0),
                font=FONT_BODY, size=Pt(14), color=RGBColor(0xB0, 0xBB, 0xD8))

    add_textbox(sl, "research@incisiv.io  |  incisiv.io",
                Inches(0.6), Inches(4.95), Inches(8), Inches(0.55),
                font=FONT_HEAD, size=Pt(18), bold=True, color=ORANGE)

    add_textbox(sl, "© 2026 Incisiv. All rights reserved.",
                Inches(0.6), H - Inches(0.7), Inches(6), Inches(0.4),
                font=FONT_BODY, size=Pt(9), color=TEXT3)


# ── Build all slides ─────────────────────────────────────────────────────────
slide_cover(prs)
slide_definition(prs)
slide_opportunity(prs)
slide_use_cases(prs)
slide_maturity(prs)
slide_benefits(prs)
slide_challenges(prs)
slide_roadahead(prs)
slide_close(prs)

out = "/home/user/Brand-Guideline/Incisiv_Agentic_AI_in_Retail_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
